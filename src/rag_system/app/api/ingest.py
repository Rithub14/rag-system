import io
import os
import uuid
from typing import Optional

from fastapi import APIRouter, File, Form, HTTPException, Request, UploadFile
from pypdf import PdfReader

from ..observability.ratelimit import rate_limiter
from ..retrieval.chunking import chunk_text
from ..retrieval.embeddings import embed_texts
from ..retrieval.vector import ensure_schema, get_client, store_chunks

router = APIRouter()


@router.post("/ingest/file")
async def ingest_file(
    request: Request,
    file: UploadFile = File(...),
    doc_id: Optional[str] = Form(None),
) -> dict:
    filename = (file.filename or "").lower()
    if not filename:
        raise HTTPException(status_code=400, detail="Missing filename.")

    browser_id = request.cookies.get("browser_id")
    session_id = request.headers.get("x-session-id")
    client_ip = request.headers.get("x-forwarded-for", request.client.host)
    effective_user_id = browser_id or session_id or client_ip
    rate_limiter.check("upload", effective_user_id, limit=1, window_seconds=3600)

    data = await file.read()
    max_mb = int(os.getenv("MAX_UPLOAD_MB", "10"))
    max_bytes = max_mb * 1024 * 1024
    if len(data) > max_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Max size is {max_mb} MB.",
        )
    text = ""

    if filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(pages).strip()
        if not text:
            try:
                from pdf2image import convert_from_bytes
                import pytesseract
            except Exception:
                raise HTTPException(
                    status_code=400,
                    detail="No extractable text found. OCR dependencies not installed.",
                )

            images = convert_from_bytes(data)
            ocr_pages = [pytesseract.image_to_string(img) or "" for img in images]
            text = "\n\n".join(ocr_pages).strip()
            if not text:
                raise HTTPException(
                    status_code=400,
                    detail="No extractable text found (OCR failed).",
                )
    elif filename.endswith(".docx"):
        try:
            import docx
        except Exception:
            raise HTTPException(status_code=400, detail="DOCX dependencies not installed.")
        doc = docx.Document(io.BytesIO(data))
        text = "\n".join(p.text for p in doc.paragraphs).strip()
    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
        except Exception:
            raise HTTPException(status_code=400, detail="PPTX dependencies not installed.")
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        text = "\n".join(parts).strip()
    elif filename.endswith(".xlsx"):
        try:
            from openpyxl import load_workbook
        except Exception:
            raise HTTPException(status_code=400, detail="XLSX dependencies not installed.")
        wb = load_workbook(io.BytesIO(data), data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join("" if v is None else str(v) for v in row)
                if row_text.strip():
                    rows.append(row_text)
        text = "\n".join(rows).strip()
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")

    if not text:
        raise HTTPException(status_code=400, detail="No extractable text found.")

    doc_id = doc_id or str(uuid.uuid4())
    metadata = {
        "user_id": effective_user_id,
        "doc_id": doc_id,
        "source": file.filename or "uploaded.pdf",
    }
    chunks = chunk_text(text, metadata)
    embeddings = embed_texts([c.content for c in chunks])

    client = get_client()
    try:
        ensure_schema(client)
        store_chunks(client, chunks, embeddings)
    finally:
        client.close()

    return {"doc_id": doc_id, "chunks": len(chunks), "filename": file.filename}
